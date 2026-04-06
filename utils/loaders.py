"""
loaders.py
Каждая функция принимает путь к файлу (str) → возвращает текст (str).
load_file() — главный роутер, вызывается из ingestor.py
"""

import json
from pathlib import Path
from typing import Optional


# ─────────────────────────────────────────────────────────────
def load_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def load_pdf(path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def load_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text.strip() for c in row.cells))
    return "\n".join(parts)


def load_xlsx(path: str) -> str:
    import pandas as pd
    xl = pd.ExcelFile(path)
    parts = []
    for sheet in xl.sheet_names:
        df = xl.parse(sheet)
        parts.append(f"=== Sheet: {sheet} ===")
        parts.append(df.to_string(index=False))
    return "\n".join(parts)


def load_csv(path: str) -> str:
    import pandas as pd
    df = pd.read_csv(path)
    return df.to_string(index=False)


def load_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def load_json(path: str) -> str:
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)
    return json.dumps(data, ensure_ascii=False, indent=2)


def load_yaml(path: str) -> str:
    import yaml
    with open(path, "r", encoding="utf-8") as f:
        data = yaml.safe_load(f)
    return yaml.dump(data, allow_unicode=True)


def load_html(path: str) -> str:
    from bs4 import BeautifulSoup
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n")


def load_image(path: str) -> str:
    """Описывает изображение через Gemini Vision → возвращает текст описания"""
    import PIL.Image
    from google import genai
    from utils.config import GEMINI_API_KEY, CHAT_MODEL

    client = genai.Client(api_key=GEMINI_API_KEY)
    img = PIL.Image.open(path)
    resp = client.models.generate_content(
        model=CHAT_MODEL,
        contents=["Опиши подробно содержимое этого изображения для поиска.", img],
    )
    return resp.text


# ── Роутер: расширение → функция ─────────────────────────────
LOADERS: dict = {
    ".txt":  load_txt,
    ".md":   load_txt,
    ".py":   load_txt,
    ".js":   load_txt,
    ".ts":   load_txt,
    ".pdf":  load_pdf,
    ".docx": load_docx,
    ".xlsx": load_xlsx,
    ".xls":  load_xlsx,
    ".csv":  load_csv,
    ".pptx": load_pptx,
    ".json": load_json,
    ".yaml": load_yaml,
    ".yml":  load_yaml,
    ".html": load_html,
    ".htm":  load_html,
    ".png":  load_image,
    ".jpg":  load_image,
    ".jpeg": load_image,
    ".webp": load_image,
}


def load_file(path: str) -> Optional[str]:
    """
    Главная функция. Принимает путь к файлу любого типа.
    Возвращает текст или None если файл нечитаемый.
    """
    ext = Path(path).suffix.lower()
    loader = LOADERS.get(ext, load_txt)   # fallback → txt
    try:
        text = loader(path)
        return text if text and text.strip() else None
    except Exception as e:
        raise ValueError(f"Не удалось загрузить файл {Path(path).name}: {e}")
